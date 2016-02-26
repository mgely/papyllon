import subprocess
import win32gui
import win32con
import win32com
import time
import ctypes
import os 
import json
from Tkinter import *
import ttk
from pptx import Presentation
from pptx.util import Inches, Pt

def byteify(input):
    '''Utility function to parse all the unicode expressions in a dictionary
    to a string
    '''
    if isinstance(input, dict):
        return {byteify(key):byteify(value) for key,value in input.iteritems()}
    elif isinstance(input, list):
        return [byteify(element) for element in input]
    elif isinstance(input, unicode):
        return input.encode('utf-8')
    else:
        return input

def create_png(data_folder):

    #############################################
    # Load spyview exe firectory
    #############################################

    
    global_settings = os.path.join(os.path.dirname(__file__), "global.json")
    with open(global_settings,"r") as f:
        global_setup =  byteify(json.load(f))
    spyview_path = global_setup["spyview_console_path"]


    #############################################
    # Find data file
    #############################################
    for file in os.listdir(data_folder):
        if file.endswith(".dat"):
            data_folder_name = file
            data_directory_path = os.path.join(data_folder, data_folder_name)


    #############################################
    # Open spyview in subprocess
    #############################################
    SW_MINIMIZE = 6
    info = subprocess.STARTUPINFO()
    info.dwFlags = subprocess.STARTF_USESHOWWINDOW
    info.wShowWindow = SW_MINIMIZE
    p = subprocess.Popen([spyview_path
        ,data_directory_path]
        , startupinfo=info)



    #############################################
    # Go through all open windows until the
    # open spyview window is found
    #############################################
    EnumWindows = ctypes.windll.user32.EnumWindows
    EnumWindowsProc = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.POINTER(ctypes.c_int), ctypes.POINTER(ctypes.c_int))
    GetWindowText = ctypes.windll.user32.GetWindowTextW
    GetWindowTextLength = ctypes.windll.user32.GetWindowTextLengthW
    IsWindowVisible = ctypes.windll.user32.IsWindowVisible

    spyview_started = False
    while spyview_started == False:
        time.sleep(0.1)
        titles = []
        def foreach_window(hwnd, lParam):
            if IsWindowVisible(hwnd):
                length = GetWindowTextLength(hwnd)
                buff = ctypes.create_unicode_buffer(length + 1)
                GetWindowText(hwnd, buff, length + 1)
                titles.append(buff.value)
            return True
        EnumWindows(EnumWindowsProc(foreach_window), 0)


        titles = [element.encode('utf-8') for element in titles]

        for title in titles:
            if data_folder_name in title:
                spyview_started = True
                spyview_title  = title

    time.sleep(0.3)



    #############################################
    # Go through all open windows to see
    # if the loading file is still open
    # In which case.. close it with alt+tab, alt+F4 hack
    #############################################
    import win32com.client
    shell = win32com.client.Dispatch("WScript.Shell")


    titles = []
    def foreach_window(hwnd, lParam):
        if IsWindowVisible(hwnd):
            length = GetWindowTextLength(hwnd)
            buff = ctypes.create_unicode_buffer(length + 1)
            GetWindowText(hwnd, buff, length + 1)
            titles.append(buff.value)
        return True
    EnumWindows(EnumWindowsProc(foreach_window), 0)


    titles = [element.encode('utf-8') for element in titles]

    for title in titles:
        if 'Loading file...' in title:
            hwnd_loader = win32gui.FindWindow(None, title)
            ctypes.windll.User32.SetForegroundWindow(hwnd_loader)
            # shell.SendKeys("%+{TAB}")
            # shell.SendKeys("%{TAB}")
            # shell.SendKeys("%{TAB}")
            shell.SendKeys("%{F4}")
            # shell.SendKeys("%{P}")
            # hwnd = win32gui.FindWindow('spyview.exe', spyview_title)
            # ctypes.windll.User32.ShowWindow(hwnd, win32con.SW_MINIMIZE)
            time.sleep(0.1)
            


    #############################################
    # Send Alt+P to spyview then close
    #############################################

    hwnd = win32gui.FindWindow('spyview_console.exe', spyview_title)
    ctypes.windll.User32.ShowWindow(hwnd, win32con.SW_MINIMIZE)
    ctypes.windll.User32.SetForegroundWindow(hwnd)
    time.sleep(0.1)
    shell.SendKeys("%{P}")


    time_to_wait = os.path.getsize(data_directory_path) / 1.0e8
    if time_to_wait < 1:
        time_to_wait = 1

    time.sleep(time_to_wait)
    #hex codes: https://msdn.microsoft.com/en-us/library/dd375731(v=VS.85).aspx
    win32gui.SendMessage(hwnd, 256, 0x51, 0)
    win32gui.SendMessage(hwnd, 257, 0x51, 65539)

def add_measurement_to_ppt(ppt_file,data_directory):

    print "adding "+data_directory

    meta_information = False
    for file in os.listdir(data_directory):
        if file.endswith(".png") or file.endswith(".PNG"):
            png = os.path.join(data_directory, file)
        if file.endswith(".json"):
            if 'timing' in file:
                with open(os.path.join(data_directory, file),"r") as f:
                    timing = byteify(json.load(f))
                    meta_information = True
            elif 'setup' in file:
                with open(os.path.join(data_directory, file),"r") as f:
                    setup = byteify(json.load(f))
            elif 'settings' in file:
                with open(os.path.join(data_directory, file),"r") as f:
                    settings = byteify(json.load(f))
                


    # load a presentation
    file_already_open = True
    while file_already_open == True:
        try: 
            with open(ppt_file,'a') as f:
                pass
            file_already_open = False
        except IOError, e:
            print str(e)

            root = Tk()
            root.title("Error opening file")
            root.lift()
            mainframe = ttk.Frame(root,padding = (3,3,3,3))
            mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

            mainframe.columnconfigure(0,weight = 1)
            mainframe.rowconfigure(0,weight = 1)

            ttk.Label(mainframe,
                text = "Please close the raw data ppt: "+ppt_file
                ).grid(column = 0, row = 0, sticky = E)

            ttk.Button(mainframe,
                text = "Done",
                command = root.destroy
                ).grid(column = 0, row = 1)
            root.mainloop()


    prs = Presentation(ppt_file)

    # add measurement info
    if meta_information == True:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.text = settings["measurement_name"] + ": " + setup["experiment"] + "\n"
        p.font.size = Pt(18)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = "started: " + timing["start"] + "\n" +\
                "ended: " + timing["end"] + "\n" +\
                "lasted: " + timing["duration"] + "\n" 
        p.font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = "Notes: " + setup["notes"] + "\n" 
        p.font.size = Pt(16)

        p = tf.add_paragraph()
        p.text = "Settings: " + "\n" + json.dumps(settings, indent=4)
        p.font.size = Pt(9)


    # add png
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(png, 0, 0,Pt(720), Pt(540))

    # save it
    prs.save(ppt_file)

def ask_user_to_generate_png(data_folder):

    #############################################
    # Load spyview exe firectory
    #############################################

    
    global_settings = os.path.join(os.path.dirname(__file__), "global.json")
    with open(global_settings,"r") as f:
        global_setup =  byteify(json.load(f))
    spyview_path = global_setup["spyview_console_path"]


    #############################################
    # Find data file
    #############################################
    for file in os.listdir(data_folder):
        if file.endswith(".dat"):
            data_folder_name = file
            data_directory_path = os.path.join(data_folder, data_folder_name)


    #############################################
    # Open spyview in subprocess
    #############################################
    SW_MINIMIZE = 6
    info = subprocess.STARTUPINFO()
    info.dwFlags = subprocess.STARTF_USESHOWWINDOW
    info.wShowWindow = SW_MINIMIZE
    p = subprocess.Popen([spyview_path
        ,data_directory_path]
        , startupinfo=info)



    #############################################
    # Go through all open windows until the
    # open spyview window is found
    #############################################
    EnumWindows = ctypes.windll.user32.EnumWindows
    EnumWindowsProc = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.POINTER(ctypes.c_int), ctypes.POINTER(ctypes.c_int))
    GetWindowText = ctypes.windll.user32.GetWindowTextW
    GetWindowTextLength = ctypes.windll.user32.GetWindowTextLengthW
    IsWindowVisible = ctypes.windll.user32.IsWindowVisible

    spyview_started = False
    while spyview_started == False:
        time.sleep(0.1)
        titles = []
        def foreach_window(hwnd, lParam):
            if IsWindowVisible(hwnd):
                length = GetWindowTextLength(hwnd)
                buff = ctypes.create_unicode_buffer(length + 1)
                GetWindowText(hwnd, buff, length + 1)
                titles.append(buff.value)
            return True
        EnumWindows(EnumWindowsProc(foreach_window), 0)


        titles = [element.encode('utf-8') for element in titles]

        for title in titles:
            if data_folder_name in title:
                spyview_started = True
                spyview_title  = title

    time.sleep(0.3)

    #############################################
    # Show spyview
    ############################################# 
    hwnd = win32gui.FindWindow('spyview_console.exe', spyview_title)
    ctypes.windll.User32.ShowWindow(hwnd, win32con.SW_MINIMIZE)
    ctypes.windll.User32.SetForegroundWindow(hwnd)


    #############################################
    # Ask user to create png
    #############################################   
    root = Tk()
    root.title("Waiting for png")
    root.lift()
    mainframe = ttk.Frame(root,padding = (3,3,3,3))
    mainframe.grid(column = 0, row = 0, sticky = (N,W,E,S))

    mainframe.columnconfigure(0,weight = 1)
    mainframe.rowconfigure(0,weight = 1)

    ttk.Label(mainframe,
        text = "Please generate a png file for raw data powerpoint..."
        ).grid(column = 0, row = 0, sticky = E)

    ttk.Button(mainframe,
        text = "Done",
        command = root.destroy
        ).grid(column = 0, row = 1)
    root.mainloop()

